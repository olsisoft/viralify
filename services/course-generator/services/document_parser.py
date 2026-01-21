"""
Document Parser Service

Extracts text content from various document formats.
Supports: PDF, DOCX, DOC, PPTX, PPT, TXT, MD, XLSX, XLS, CSV
"""
import csv
import io
import re
from pathlib import Path
from typing import List, Optional, Tuple

import fitz  # PyMuPDF for PDF
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches

from models.document_models import (
    Document,
    DocumentChunk,
    DocumentType,
    ExtractedImage,
)


class DocumentParser:
    """
    Multi-format document parser for text extraction.

    Extracts text content while preserving structure information
    like page numbers, sections, and headings.
    """

    # Chunk configuration
    DEFAULT_CHUNK_SIZE = 1000  # characters
    DEFAULT_CHUNK_OVERLAP = 200  # characters

    def __init__(
        self,
        chunk_size: int = DEFAULT_CHUNK_SIZE,
        chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    async def parse_document(
        self,
        content: bytes,
        filename: str,
        document_type: DocumentType,
    ) -> Tuple[str, List[DocumentChunk], dict]:
        """
        Parse document and extract text content.

        Args:
            content: Raw file bytes
            filename: Original filename
            document_type: Type of document

        Returns:
            Tuple of (raw_text, chunks, metadata)
        """
        print(f"[PARSER] Parsing document: {filename} (type: {document_type})", flush=True)

        # Route to appropriate parser
        parser_map = {
            DocumentType.PDF: self._parse_pdf,
            DocumentType.DOCX: self._parse_docx,
            DocumentType.DOC: self._parse_doc,
            DocumentType.PPTX: self._parse_pptx,
            DocumentType.PPT: self._parse_ppt,
            DocumentType.TXT: self._parse_txt,
            DocumentType.MD: self._parse_markdown,
            DocumentType.XLSX: self._parse_xlsx,
            DocumentType.XLS: self._parse_xls,
            DocumentType.CSV: self._parse_csv,
        }

        parser = parser_map.get(document_type)
        if not parser:
            raise ValueError(f"Unsupported document type: {document_type}")

        # Extract raw text and metadata
        raw_text, metadata = await parser(content, filename)

        # Create chunks
        chunks = self._create_chunks(raw_text, metadata)

        print(f"[PARSER] Extracted {len(raw_text)} chars, {len(chunks)} chunks", flush=True)

        return raw_text, chunks, metadata

    async def _parse_pdf(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse PDF document using PyMuPDF"""
        metadata = {
            "pages": [],
            "page_count": 0,
            "has_images": False,
            "has_tables": False,
        }

        text_parts = []

        try:
            doc = fitz.open(stream=content, filetype="pdf")
            metadata["page_count"] = len(doc)

            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text("text")

                # Track page boundaries
                metadata["pages"].append({
                    "number": page_num,
                    "char_start": len("".join(text_parts)),
                    "char_end": len("".join(text_parts)) + len(page_text),
                })

                text_parts.append(page_text)

                # Check for images
                if page.get_images():
                    metadata["has_images"] = True

                # Check for tables (heuristic: multiple aligned columns)
                if self._detect_table_structure(page_text):
                    metadata["has_tables"] = True

            doc.close()

        except Exception as e:
            print(f"[PARSER] PDF parsing error: {e}", flush=True)
            raise ValueError(f"Failed to parse PDF: {e}")

        return "\n\n".join(text_parts), metadata

    async def _parse_docx(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse DOCX document"""
        metadata = {
            "paragraphs": 0,
            "headings": [],
            "has_tables": False,
            "has_images": False,
        }

        text_parts = []

        try:
            doc = DocxDocument(io.BytesIO(content))

            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    text_parts.append(text)
                    metadata["paragraphs"] += 1

                    # Detect headings
                    if para.style and para.style.name.startswith('Heading'):
                        metadata["headings"].append({
                            "level": int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1,
                            "text": text,
                            "char_pos": len("\n\n".join(text_parts[:-1])) if len(text_parts) > 1 else 0,
                        })

            # Check for tables
            if doc.tables:
                metadata["has_tables"] = True
                for table in doc.tables:
                    table_text = self._extract_table_text(table)
                    if table_text:
                        text_parts.append(table_text)

            # Check for images
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    metadata["has_images"] = True
                    break

        except Exception as e:
            print(f"[PARSER] DOCX parsing error: {e}", flush=True)
            raise ValueError(f"Failed to parse DOCX: {e}")

        return "\n\n".join(text_parts), metadata

    async def _parse_doc(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse legacy DOC document (limited support)"""
        # For legacy .doc, we try to extract text using basic patterns
        # Full support would require antiword or similar
        metadata = {"format": "legacy_doc", "limited_support": True}

        try:
            # Try to decode as text (works for some .doc files)
            text = content.decode('utf-8', errors='ignore')

            # Clean up binary artifacts
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', ' ', text)
            text = re.sub(r'\s+', ' ', text)

            # Extract readable portions
            readable_parts = re.findall(r'[A-Za-z0-9\s,.!?;:\'"()-]{20,}', text)

            if readable_parts:
                return "\n".join(readable_parts), metadata
            else:
                raise ValueError("Could not extract readable text from DOC file")

        except Exception as e:
            print(f"[PARSER] DOC parsing error: {e}", flush=True)
            raise ValueError(f"Failed to parse DOC: {e}. Consider converting to DOCX.")

    async def _parse_pptx(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse PowerPoint PPTX document"""
        metadata = {
            "slide_count": 0,
            "slides": [],
            "has_images": False,
            "has_charts": False,
        }

        text_parts = []

        try:
            prs = Presentation(io.BytesIO(content))
            metadata["slide_count"] = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                slide_title = None

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_texts.append(text)

                        # Detect title
                        if shape.is_placeholder and hasattr(shape, "placeholder_format"):
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                slide_title = text

                    # Check for images
                    if shape.shape_type == 13:  # Picture
                        metadata["has_images"] = True

                    # Check for charts
                    if shape.has_chart:
                        metadata["has_charts"] = True

                if slide_texts:
                    slide_content = "\n".join(slide_texts)
                    text_parts.append(f"[Slide {slide_num}]\n{slide_content}")

                    metadata["slides"].append({
                        "number": slide_num,
                        "title": slide_title,
                        "char_start": len("\n\n".join(text_parts[:-1])) if len(text_parts) > 1 else 0,
                    })

        except Exception as e:
            print(f"[PARSER] PPTX parsing error: {e}", flush=True)
            raise ValueError(f"Failed to parse PPTX: {e}")

        return "\n\n".join(text_parts), metadata

    async def _parse_ppt(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse legacy PPT document (limited support)"""
        metadata = {"format": "legacy_ppt", "limited_support": True}

        # Similar approach to DOC - try text extraction
        try:
            text = content.decode('utf-8', errors='ignore')
            text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', ' ', text)
            readable_parts = re.findall(r'[A-Za-z0-9\s,.!?;:\'"()-]{20,}', text)

            if readable_parts:
                return "\n".join(readable_parts), metadata
            else:
                raise ValueError("Could not extract readable text from PPT file")

        except Exception as e:
            print(f"[PARSER] PPT parsing error: {e}", flush=True)
            raise ValueError(f"Failed to parse PPT: {e}. Consider converting to PPTX.")

    async def _parse_txt(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse plain text file"""
        metadata = {"format": "plain_text"}

        # Try different encodings
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']

        for encoding in encodings:
            try:
                text = content.decode(encoding)
                metadata["encoding"] = encoding
                metadata["line_count"] = text.count('\n') + 1
                metadata["word_count"] = len(text.split())
                return text, metadata
            except UnicodeDecodeError:
                continue

        raise ValueError("Could not decode text file with supported encodings")

    async def _parse_markdown(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse Markdown file"""
        text, metadata = await self._parse_txt(content, filename)
        metadata["format"] = "markdown"

        # Extract headings
        headings = []
        for match in re.finditer(r'^(#{1,6})\s+(.+)$', text, re.MULTILINE):
            headings.append({
                "level": len(match.group(1)),
                "text": match.group(2),
                "char_pos": match.start(),
            })

        metadata["headings"] = headings

        # Detect code blocks
        code_blocks = re.findall(r'```[\s\S]*?```', text)
        metadata["code_block_count"] = len(code_blocks)

        return text, metadata

    async def _parse_xlsx(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse Excel XLSX file"""
        metadata = {
            "sheet_count": 0,
            "sheets": [],
            "total_rows": 0,
        }

        text_parts = []

        try:
            wb = load_workbook(io.BytesIO(content), data_only=True)
            metadata["sheet_count"] = len(wb.sheetnames)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = []
                row_count = 0

                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))

                    if row_values:
                        sheet_text.append(" | ".join(row_values))
                        row_count += 1

                if sheet_text:
                    text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_text))
                    metadata["sheets"].append({
                        "name": sheet_name,
                        "rows": row_count,
                    })
                    metadata["total_rows"] += row_count

            wb.close()

        except Exception as e:
            print(f"[PARSER] XLSX parsing error: {e}", flush=True)
            raise ValueError(f"Failed to parse XLSX: {e}")

        return "\n\n".join(text_parts), metadata

    async def _parse_xls(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse legacy XLS file"""
        # For legacy .xls, we'd need xlrd library
        # For now, provide limited support message
        metadata = {"format": "legacy_xls", "limited_support": True}

        try:
            # Try to import xlrd if available
            import xlrd

            wb = xlrd.open_workbook(file_contents=content)
            text_parts = []

            for sheet in wb.sheets():
                sheet_text = []
                for row_idx in range(sheet.nrows):
                    row_values = [str(cell.value) for cell in sheet.row(row_idx) if cell.value]
                    if row_values:
                        sheet_text.append(" | ".join(row_values))

                if sheet_text:
                    text_parts.append(f"[Sheet: {sheet.name}]\n" + "\n".join(sheet_text))

            return "\n\n".join(text_parts), metadata

        except ImportError:
            raise ValueError("XLS format requires xlrd library. Consider converting to XLSX.")
        except Exception as e:
            print(f"[PARSER] XLS parsing error: {e}", flush=True)
            raise ValueError(f"Failed to parse XLS: {e}")

    async def _parse_csv(self, content: bytes, filename: str) -> Tuple[str, dict]:
        """Parse CSV file"""
        metadata = {
            "format": "csv",
            "row_count": 0,
            "column_count": 0,
            "headers": [],
        }

        text_parts = []

        try:
            # Detect encoding
            text = content.decode('utf-8')
        except UnicodeDecodeError:
            text = content.decode('latin-1')

        try:
            # Detect delimiter
            sniffer = csv.Sniffer()
            dialect = sniffer.sniff(text[:1024])

            reader = csv.reader(io.StringIO(text), dialect)
            rows = list(reader)

            if rows:
                # First row as headers
                metadata["headers"] = rows[0]
                metadata["column_count"] = len(rows[0])
                metadata["row_count"] = len(rows)

                # Convert to text
                for row in rows:
                    text_parts.append(" | ".join(row))

        except csv.Error:
            # Fallback: treat as pipe-delimited or tab-delimited
            lines = text.strip().split('\n')
            for line in lines:
                text_parts.append(line)
            metadata["row_count"] = len(lines)

        return "\n".join(text_parts), metadata

    def _create_chunks(
        self,
        text: str,
        metadata: dict,
    ) -> List[DocumentChunk]:
        """
        Split text into overlapping chunks for embedding.

        Uses semantic boundaries (paragraphs, sentences) when possible.
        """
        chunks = []

        # Clean text
        text = text.strip()
        if not text:
            return chunks

        # Split by paragraphs first
        paragraphs = re.split(r'\n\s*\n', text)

        current_chunk = ""
        current_start = 0
        chunk_index = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # If adding this paragraph exceeds chunk size
            if len(current_chunk) + len(para) + 2 > self.chunk_size:
                # Save current chunk if it has content
                if current_chunk:
                    chunk = DocumentChunk(
                        document_id="",  # Will be set when saved
                        content=current_chunk,
                        chunk_index=chunk_index,
                        page_number=self._get_page_for_position(current_start, metadata),
                        section_title=self._get_section_for_position(current_start, metadata),
                        token_count=self._estimate_tokens(current_chunk),
                    )
                    chunks.append(chunk)
                    chunk_index += 1

                    # Start new chunk with overlap
                    overlap_text = self._get_overlap_text(current_chunk)
                    current_chunk = overlap_text + "\n\n" + para if overlap_text else para
                    current_start = current_start + len(current_chunk) - len(overlap_text) - len(para) - 2
                else:
                    current_chunk = para
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para

        # Don't forget the last chunk
        if current_chunk:
            chunk = DocumentChunk(
                document_id="",
                content=current_chunk,
                chunk_index=chunk_index,
                page_number=self._get_page_for_position(current_start, metadata),
                section_title=self._get_section_for_position(current_start, metadata),
                token_count=self._estimate_tokens(current_chunk),
            )
            chunks.append(chunk)

        return chunks

    def _get_overlap_text(self, text: str) -> str:
        """Get overlap text from end of chunk"""
        if len(text) <= self.chunk_overlap:
            return text

        # Try to find sentence boundary
        overlap_text = text[-self.chunk_overlap:]
        sentence_start = overlap_text.find('. ')

        if sentence_start > 0:
            return overlap_text[sentence_start + 2:]

        return overlap_text

    def _get_page_for_position(self, char_pos: int, metadata: dict) -> Optional[int]:
        """Get page number for character position"""
        pages = metadata.get("pages", [])
        for page in pages:
            if page.get("char_start", 0) <= char_pos < page.get("char_end", float('inf')):
                return page.get("number")
        return None

    def _get_section_for_position(self, char_pos: int, metadata: dict) -> Optional[str]:
        """Get section title for character position"""
        headings = metadata.get("headings", [])
        current_heading = None

        for heading in headings:
            if heading.get("char_pos", 0) <= char_pos:
                current_heading = heading.get("text")
            else:
                break

        return current_heading

    def _estimate_tokens(self, text: str) -> int:
        """Estimate token count (rough approximation)"""
        # Average ~4 characters per token for English
        return len(text) // 4

    def _detect_table_structure(self, text: str) -> bool:
        """Detect if text contains table-like structure"""
        lines = text.split('\n')
        aligned_lines = 0

        for line in lines:
            # Check for multiple aligned columns (spaces or tabs)
            if re.search(r'\s{2,}[\w\d].*\s{2,}[\w\d]', line):
                aligned_lines += 1

        return aligned_lines >= 3

    def _extract_table_text(self, table) -> str:
        """Extract text from DOCX table"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        return "\n".join(rows) if rows else ""

    async def extract_images(
        self,
        content: bytes,
        filename: str,
        document_type: DocumentType,
        document_id: str,
        output_dir: str,
    ) -> List[ExtractedImage]:
        """
        Extract images from document.

        Args:
            content: Raw file bytes
            filename: Original filename
            document_type: Type of document
            document_id: Document ID for references
            output_dir: Directory to save extracted images

        Returns:
            List of ExtractedImage objects
        """
        print(f"[PARSER] Extracting images from: {filename} (type: {document_type})", flush=True)

        # Create output directory
        import os
        os.makedirs(output_dir, exist_ok=True)

        images = []

        if document_type == DocumentType.PDF:
            images = await self._extract_pdf_images(content, document_id, output_dir)
        elif document_type == DocumentType.PPTX:
            images = await self._extract_pptx_images(content, document_id, output_dir)
        elif document_type == DocumentType.DOCX:
            images = await self._extract_docx_images(content, document_id, output_dir)

        print(f"[PARSER] Extracted {len(images)} images from {filename}", flush=True)
        return images

    async def _extract_pdf_images(
        self,
        content: bytes,
        document_id: str,
        output_dir: str,
    ) -> List[ExtractedImage]:
        """Extract images from PDF using PyMuPDF"""
        import os
        from PIL import Image as PILImage

        images = []

        try:
            doc = fitz.open(stream=content, filetype="pdf")

            for page_num, page in enumerate(doc, 1):
                # Get images on this page
                image_list = page.get_images(full=True)

                # Get text around images for context
                page_text = page.get_text("text")[:500]  # First 500 chars as context

                for img_index, img in enumerate(image_list):
                    try:
                        # Extract image
                        xref = img[0]
                        base_image = doc.extract_image(xref)

                        if base_image:
                            image_bytes = base_image["image"]
                            image_ext = base_image.get("ext", "png")

                            # Check if image is large enough to be useful (min 100x100)
                            pil_img = PILImage.open(io.BytesIO(image_bytes))
                            width, height = pil_img.size

                            if width >= 100 and height >= 100:
                                # Save image
                                img_filename = f"{document_id}_p{page_num}_img{img_index}.{image_ext}"
                                img_path = os.path.join(output_dir, img_filename)

                                with open(img_path, "wb") as f:
                                    f.write(image_bytes)

                                # Detect image type based on dimensions and aspect ratio
                                detected_type = self._detect_image_type(width, height, image_ext)

                                extracted_img = ExtractedImage(
                                    document_id=document_id,
                                    file_path=img_path,
                                    file_name=img_filename,
                                    image_format=image_ext,
                                    width=width,
                                    height=height,
                                    file_size_bytes=len(image_bytes),
                                    page_number=page_num,
                                    context_text=page_text,
                                    detected_type=detected_type,
                                )
                                images.append(extracted_img)

                    except Exception as img_error:
                        print(f"[PARSER] Error extracting image {img_index} from page {page_num}: {img_error}", flush=True)
                        continue

            doc.close()

        except Exception as e:
            print(f"[PARSER] PDF image extraction error: {e}", flush=True)

        return images

    async def _extract_pptx_images(
        self,
        content: bytes,
        document_id: str,
        output_dir: str,
    ) -> List[ExtractedImage]:
        """Extract images from PowerPoint"""
        import os
        from PIL import Image as PILImage

        images = []

        try:
            prs = Presentation(io.BytesIO(content))

            for slide_num, slide in enumerate(prs.slides, 1):
                # Get slide text for context
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                slide_text = slide_text[:500]

                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        if shape.shape_type == 13:  # Picture
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = image.ext

                            # Get dimensions
                            pil_img = PILImage.open(io.BytesIO(image_bytes))
                            width, height = pil_img.size

                            if width >= 100 and height >= 100:
                                img_filename = f"{document_id}_s{slide_num}_img{shape_idx}.{image_ext}"
                                img_path = os.path.join(output_dir, img_filename)

                                with open(img_path, "wb") as f:
                                    f.write(image_bytes)

                                detected_type = self._detect_image_type(width, height, image_ext)

                                extracted_img = ExtractedImage(
                                    document_id=document_id,
                                    file_path=img_path,
                                    file_name=img_filename,
                                    image_format=image_ext,
                                    width=width,
                                    height=height,
                                    file_size_bytes=len(image_bytes),
                                    page_number=slide_num,
                                    context_text=slide_text,
                                    detected_type=detected_type,
                                )
                                images.append(extracted_img)

                    except Exception as shape_error:
                        continue

        except Exception as e:
            print(f"[PARSER] PPTX image extraction error: {e}", flush=True)

        return images

    async def _extract_docx_images(
        self,
        content: bytes,
        document_id: str,
        output_dir: str,
    ) -> List[ExtractedImage]:
        """Extract images from Word document"""
        import os
        from PIL import Image as PILImage

        images = []

        try:
            doc = DocxDocument(io.BytesIO(content))

            # Get all images from relationships
            img_index = 0
            for rel_id, rel in doc.part.rels.items():
                if "image" in rel.reltype:
                    try:
                        image_part = rel.target_part
                        image_bytes = image_part.blob
                        content_type = image_part.content_type

                        # Determine extension from content type
                        ext_map = {
                            "image/png": "png",
                            "image/jpeg": "jpg",
                            "image/gif": "gif",
                            "image/bmp": "bmp",
                        }
                        image_ext = ext_map.get(content_type, "png")

                        # Get dimensions
                        pil_img = PILImage.open(io.BytesIO(image_bytes))
                        width, height = pil_img.size

                        if width >= 100 and height >= 100:
                            img_filename = f"{document_id}_img{img_index}.{image_ext}"
                            img_path = os.path.join(output_dir, img_filename)

                            with open(img_path, "wb") as f:
                                f.write(image_bytes)

                            detected_type = self._detect_image_type(width, height, image_ext)

                            extracted_img = ExtractedImage(
                                document_id=document_id,
                                file_path=img_path,
                                file_name=img_filename,
                                image_format=image_ext,
                                width=width,
                                height=height,
                                file_size_bytes=len(image_bytes),
                                detected_type=detected_type,
                            )
                            images.append(extracted_img)
                            img_index += 1

                    except Exception as img_error:
                        continue

        except Exception as e:
            print(f"[PARSER] DOCX image extraction error: {e}", flush=True)

        return images

    def _detect_image_type(self, width: int, height: int, ext: str) -> str:
        """Detect the type of image based on characteristics"""
        aspect_ratio = width / height if height > 0 else 1

        # Screenshots tend to be wider
        if aspect_ratio > 1.5 and width > 800:
            return "screenshot"

        # Charts and diagrams are often square-ish
        if 0.7 <= aspect_ratio <= 1.4:
            if width > 400:
                return "diagram"
            return "icon"

        # Photos in presentations are often landscape
        if ext in ["jpg", "jpeg"] and aspect_ratio > 1.2:
            return "photo"

        # Default to diagram for educational content
        return "diagram"


# URL and YouTube parsers (separate from file parsing)
class WebContentParser:
    """Parser for web-based content sources"""

    async def parse_url(self, url: str) -> Tuple[str, dict]:
        """
        Fetch and parse content from URL.

        Note: Requires httpx and beautifulsoup4
        """
        import httpx
        from bs4 import BeautifulSoup

        metadata = {
            "source_url": url,
            "format": "web_page",
        }

        try:
            # Use browser-like headers to avoid 403 blocks
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8",
                "Accept-Language": "en-US,en;q=0.9,fr;q=0.8",
                "Accept-Encoding": "gzip, deflate, br",
                "Connection": "keep-alive",
                "Upgrade-Insecure-Requests": "1",
                "Sec-Fetch-Dest": "document",
                "Sec-Fetch-Mode": "navigate",
                "Sec-Fetch-Site": "none",
                "Sec-Fetch-User": "?1",
                "Cache-Control": "max-age=0",
            }
            async with httpx.AsyncClient(timeout=30.0, headers=headers) as client:
                response = await client.get(url, follow_redirects=True)
                response.raise_for_status()

                metadata["status_code"] = response.status_code
                metadata["content_type"] = response.headers.get("content-type", "")

                # Parse HTML
                soup = BeautifulSoup(response.text, 'html.parser')

                # Remove script and style elements
                for element in soup(['script', 'style', 'nav', 'footer', 'header']):
                    element.decompose()

                # Extract title
                title = soup.find('title')
                metadata["title"] = title.get_text().strip() if title else None

                # Extract main content
                main_content = soup.find('main') or soup.find('article') or soup.find('body')

                if main_content:
                    text = main_content.get_text(separator='\n', strip=True)
                else:
                    text = soup.get_text(separator='\n', strip=True)

                # Clean up whitespace
                text = re.sub(r'\n\s*\n', '\n\n', text)

                return text, metadata

        except httpx.HTTPStatusError as e:
            print(f"[PARSER] URL HTTP error: {e.response.status_code}", flush=True)
            if e.response.status_code == 404:
                raise ValueError(f"Page non trouvée (404). Vérifiez que l'URL est correcte: {url}")
            elif e.response.status_code == 403:
                raise ValueError(f"Accès refusé (403). Cette page bloque les robots.")
            elif e.response.status_code >= 500:
                raise ValueError(f"Erreur serveur ({e.response.status_code}). Réessayez plus tard.")
            else:
                raise ValueError(f"Erreur HTTP {e.response.status_code} lors de la récupération de l'URL")
        except httpx.TimeoutException:
            print(f"[PARSER] URL timeout: {url}", flush=True)
            raise ValueError(f"Timeout: La page met trop de temps à répondre")
        except Exception as e:
            print(f"[PARSER] URL parsing error: {e}", flush=True)
            raise ValueError(f"Erreur lors de la récupération de l'URL: {str(e)[:100]}")

    async def parse_youtube(self, url: str) -> Tuple[str, dict]:
        """
        Extract transcript from YouTube video.

        Strategy: Use yt-dlp first (more reliable), fallback to youtube_transcript_api.
        """
        import re

        metadata = {
            "source_url": url,
            "format": "youtube_transcript",
        }

        # Extract video ID - support multiple URL formats
        video_id_match = re.search(
            r'(?:youtube\.com/watch\?v=|youtube\.com/watch\?.*v=|youtu\.be/|youtube\.com/embed/|youtube\.com/v/|youtube\.com/shorts/)([a-zA-Z0-9_-]{11})',
            url
        )

        if not video_id_match:
            video_id_match = re.search(r'[?&]v=([a-zA-Z0-9_-]{11})', url)

        if not video_id_match:
            raise ValueError(f"Format d'URL YouTube invalide: {url}")

        video_id = video_id_match.group(1)
        metadata["video_id"] = video_id
        print(f"[PARSER] Extracting transcript for YouTube video: {video_id}", flush=True)

        errors = []

        # Method 1: Try yt-dlp first (more reliable)
        try:
            print(f"[PARSER] Trying yt-dlp method...", flush=True)
            text, yt_metadata = await self._parse_youtube_with_ytdlp(url, video_id)
            metadata.update(yt_metadata)
            metadata["method"] = "yt-dlp"
            print(f"[PARSER] yt-dlp success: {len(text)} characters", flush=True)
            return text, metadata
        except Exception as e:
            errors.append(f"yt-dlp: {e}")
            print(f"[PARSER] yt-dlp failed: {e}", flush=True)

        # Method 2: Fallback to youtube_transcript_api
        try:
            print(f"[PARSER] Trying youtube_transcript_api method...", flush=True)
            text, api_metadata = await self._parse_youtube_with_api(video_id)
            metadata.update(api_metadata)
            metadata["method"] = "youtube_transcript_api"
            print(f"[PARSER] youtube_transcript_api success: {len(text)} characters", flush=True)
            return text, metadata
        except Exception as e:
            errors.append(f"transcript_api: {e}")
            print(f"[PARSER] youtube_transcript_api failed: {e}", flush=True)

        # Both methods failed
        error_summary = "; ".join(errors)
        print(f"[PARSER] All YouTube methods failed: {error_summary}", flush=True)

        # Provide user-friendly error message
        if any("unavailable" in err.lower() or "private" in err.lower() for err in errors):
            raise ValueError("Cette vidéo YouTube n'est pas disponible ou est privée.")
        elif any("disabled" in err.lower() for err in errors):
            raise ValueError("Les sous-titres sont désactivés pour cette vidéo YouTube.")
        elif any("no subtitle" in err.lower() or "aucun sous-titre" in err.lower() for err in errors):
            raise ValueError("Cette vidéo n'a pas de sous-titres disponibles.")
        else:
            raise ValueError(f"Impossible d'extraire la transcription YouTube. Vérifiez que la vidéo est publique et a des sous-titres.")

    async def _parse_youtube_with_api(self, video_id: str) -> Tuple[str, dict]:
        """
        Extract transcript using youtube_transcript_api.
        """
        from youtube_transcript_api import YouTubeTranscriptApi
        from youtube_transcript_api._errors import (
            TranscriptsDisabled,
            NoTranscriptFound,
            VideoUnavailable,
        )
        import re

        metadata = {}

        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)

            # Collect available transcripts info
            available = []
            for t in transcript_list:
                available.append({
                    'language': t.language,
                    'language_code': t.language_code,
                    'is_generated': t.is_generated,
                })
            metadata["available_transcripts"] = available
            print(f"[PARSER] Available transcripts: {[a['language_code'] for a in available]}", flush=True)

            # Try to get transcript in preferred order
            transcript_data = None

            # Priority: fr manual > en manual > fr auto > en auto > any
            for lang in ['fr', 'en']:
                try:
                    transcript = transcript_list.find_transcript([lang])
                    transcript_data = transcript.fetch()
                    metadata["language"] = lang
                    metadata["is_generated"] = transcript.is_generated
                    break
                except NoTranscriptFound:
                    continue

            # Try auto-generated if no manual found
            if not transcript_data:
                try:
                    transcript = transcript_list.find_generated_transcript(['fr', 'en'])
                    transcript_data = transcript.fetch()
                    metadata["language"] = transcript.language_code
                    metadata["is_generated"] = True
                except NoTranscriptFound:
                    pass

            # Last resort: any transcript
            if not transcript_data:
                for t in transcript_list:
                    try:
                        transcript_data = t.fetch()
                        metadata["language"] = t.language_code
                        metadata["is_generated"] = t.is_generated
                        break
                    except Exception:
                        continue

            if not transcript_data:
                raise ValueError("Aucun sous-titre extractible")

            # Combine segments
            text_parts = [segment['text'] for segment in transcript_data]
            text = " ".join(text_parts)
            text = re.sub(r'\s+', ' ', text).strip()

            if not text or len(text) < 50:
                raise ValueError("Transcription extraite trop courte ou vide")

            metadata["duration_seconds"] = transcript_data[-1]['start'] + transcript_data[-1]['duration'] if transcript_data else 0
            metadata["word_count"] = len(text.split())

            return text, metadata

        except TranscriptsDisabled:
            raise ValueError("Sous-titres désactivés pour cette vidéo")
        except VideoUnavailable:
            raise ValueError("Vidéo non disponible ou privée")
        except Exception as e:
            raise ValueError(f"Erreur API: {str(e)[:100]}")

    async def _parse_youtube_with_ytdlp(self, url: str, video_id: str) -> Tuple[str, dict]:
        """
        Extract subtitles using yt-dlp (more reliable than youtube_transcript_api).
        """
        import subprocess
        import json
        import tempfile
        import os
        import re
        import asyncio

        metadata = {}

        with tempfile.TemporaryDirectory() as tmpdir:
            # First, get video info and available subtitles
            info_cmd = [
                'yt-dlp',
                '--dump-json',
                '--skip-download',
                '--no-warnings',
                '--quiet',
                url
            ]

            try:
                # Run in executor to not block
                loop = asyncio.get_event_loop()
                result = await loop.run_in_executor(
                    None,
                    lambda: subprocess.run(
                        info_cmd,
                        capture_output=True,
                        text=True,
                        timeout=45
                    )
                )

                if result.returncode != 0:
                    error_msg = result.stderr.strip() if result.stderr else "Unknown error"
                    if "Video unavailable" in error_msg or "Private video" in error_msg:
                        raise ValueError("Vidéo non disponible ou privée")
                    raise ValueError(f"yt-dlp info failed: {error_msg[:100]}")

                video_info = json.loads(result.stdout)
                metadata['title'] = video_info.get('title', '')
                metadata['duration_seconds'] = video_info.get('duration', 0)
                metadata['uploader'] = video_info.get('uploader', '')

                # Check for subtitles
                subtitles = video_info.get('subtitles', {})
                auto_captions = video_info.get('automatic_captions', {})

                print(f"[PARSER] yt-dlp found subtitles: {list(subtitles.keys())}, auto: {list(auto_captions.keys())[:5]}", flush=True)

                if not subtitles and not auto_captions:
                    raise ValueError("Aucun sous-titre disponible pour cette vidéo")

                # Determine best subtitle language
                sub_lang = None
                is_auto = False

                # Priority: manual subtitles first
                for lang in ['fr', 'fr-FR', 'en', 'en-US', 'en-GB']:
                    if lang in subtitles:
                        sub_lang = lang
                        break

                # Then auto-generated
                if not sub_lang:
                    for lang in ['fr', 'fr-FR', 'en', 'en-US', 'en-GB']:
                        if lang in auto_captions:
                            sub_lang = lang
                            is_auto = True
                            break

                # Any available
                if not sub_lang:
                    if subtitles:
                        sub_lang = list(subtitles.keys())[0]
                    elif auto_captions:
                        sub_lang = list(auto_captions.keys())[0]
                        is_auto = True

                if not sub_lang:
                    raise ValueError("Aucune langue de sous-titre trouvée")

                metadata['language'] = sub_lang
                metadata['is_generated'] = is_auto
                print(f"[PARSER] yt-dlp: downloading subtitles in {sub_lang} (auto={is_auto})", flush=True)

                # Download subtitles
                sub_file = os.path.join(tmpdir, 'subs')
                sub_cmd = [
                    'yt-dlp',
                    '--skip-download',
                    '--write-sub' if not is_auto else '--write-auto-sub',
                    '--sub-lang', sub_lang,
                    '--sub-format', 'vtt/srt/best',
                    '--convert-subs', 'srt',
                    '--no-warnings',
                    '-o', sub_file,
                    url
                ]

                # Also try auto-subs if manual fails
                if not is_auto:
                    sub_cmd.insert(3, '--write-auto-sub')

                result = await loop.run_in_executor(
                    None,
                    lambda: subprocess.run(
                        sub_cmd,
                        capture_output=True,
                        text=True,
                        timeout=90
                    )
                )

                # Find the downloaded subtitle file (could be .srt or .vtt)
                sub_files = [f for f in os.listdir(tmpdir) if f.endswith(('.srt', '.vtt'))]
                if not sub_files:
                    # Check if there was an error
                    if result.stderr and "no subtitle" in result.stderr.lower():
                        raise ValueError("Aucun sous-titre téléchargeable")
                    raise ValueError(f"Téléchargement des sous-titres échoué")

                sub_path = os.path.join(tmpdir, sub_files[0])
                print(f"[PARSER] yt-dlp: processing subtitle file {sub_files[0]}", flush=True)

                # Parse subtitle file
                with open(sub_path, 'r', encoding='utf-8', errors='ignore') as f:
                    sub_content = f.read()

                # Extract text (handle both SRT and VTT formats)
                text = self._parse_subtitle_content(sub_content)

                if not text or len(text) < 50:
                    raise ValueError("Sous-titres extraits trop courts ou vides")

                metadata['word_count'] = len(text.split())
                metadata['char_count'] = len(text)
                print(f"[PARSER] yt-dlp: extracted {len(text)} characters, {metadata['word_count']} words", flush=True)

                return text, metadata

            except subprocess.TimeoutExpired:
                raise ValueError("Timeout - la vidéo est peut-être trop longue")
            except json.JSONDecodeError as e:
                raise ValueError(f"Erreur de parsing des infos vidéo: {e}")
            except ValueError:
                raise
            except Exception as e:
                raise ValueError(f"Erreur yt-dlp: {str(e)[:100]}")

    def _parse_subtitle_content(self, content: str) -> str:
        """
        Parse SRT or VTT subtitle content and extract clean text.
        """
        import re

        lines = content.split('\n')
        text_parts = []
        seen_lines = set()  # Avoid duplicates

        for line in lines:
            line = line.strip()

            # Skip empty lines
            if not line:
                continue

            # Skip sequence numbers (SRT)
            if line.isdigit():
                continue

            # Skip timestamps (SRT: 00:00:00,000 --> 00:00:00,000)
            if re.match(r'\d{2}:\d{2}:\d{2}[,\.]\d{3}\s*-->', line):
                continue

            # Skip VTT header and metadata
            if line.startswith('WEBVTT') or line.startswith('Kind:') or line.startswith('Language:'):
                continue

            # Skip VTT timestamps (00:00:00.000 --> 00:00:00.000)
            if re.match(r'\d{2}:\d{2}:\d{2}\.\d{3}\s*-->', line):
                continue

            # Skip VTT cue settings (position, align, etc.)
            if re.match(r'^[\d:\.]+\s*-->\s*[\d:\.]+', line):
                continue

            # Clean up HTML/VTT tags
            line = re.sub(r'<[^>]+>', '', line)  # Remove HTML tags
            line = re.sub(r'\{[^}]+\}', '', line)  # Remove style tags
            line = re.sub(r'&nbsp;', ' ', line)
            line = re.sub(r'&amp;', '&', line)
            line = re.sub(r'&lt;', '<', line)
            line = re.sub(r'&gt;', '>', line)

            # Clean up speaker labels like "[Music]" or "(applause)"
            line = re.sub(r'\[[^\]]*\]', '', line)
            line = re.sub(r'\([^)]*\)', '', line)

            line = line.strip()

            # Skip if empty after cleaning or if duplicate
            if not line or len(line) < 2:
                continue

            # Avoid exact duplicates (common in auto-generated subs)
            if line.lower() in seen_lines:
                continue
            seen_lines.add(line.lower())

            text_parts.append(line)

        # Join and clean up
        text = ' '.join(text_parts)
        text = re.sub(r'\s+', ' ', text).strip()

        return text
